import io
import csv
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pytesseract



def handle_text(uploaded_file):
    """处理文本文件"""
    with io.TextIOWrapper(uploaded_file, encoding='utf-8') as file:
        return file.read()

def handle_pdf(uploaded_file):
    """处理PDF文件，尝试更智能地识别段落和行"""
    reader = PdfReader(uploaded_file)
    text = ''

    for page in reader.pages:
        page_text = page.extract_text() or ''
        lines = page_text.split('\n')

        for line in lines:
            # 如果行以常见的句末符号结束，则在该行后添加换行符
            if line.strip() and line.strip()[-1] in ['.', '?', '!', ':']:
                text += line.strip() + '\n'
            else:
                # 否则，将这一行添加到段落中
                text += line.strip() + ' '
        text += '\n\n'  # 在页面之间添加额外的换行
    return text

def handle_docx(uploaded_file):
    """处理DOCX文件"""
    doc = Document(uploaded_file)
    return '\n'.join(para.text for para in doc.paragraphs)

def handle_excel(uploaded_file):
    """处理XLSX和XLS文件"""
    workbook = load_workbook(filename=uploaded_file)
    text = ''
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows():
            text += '\t'.join(str(cell.value) for cell in row) + '\n'
    return text

def handle_ppt(uploaded_file):
    """处理PPTX和PPT文件"""
    presentation = Presentation(uploaded_file)
    return '\n'.join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))

def handle_image(uploaded_file):
    """处理图片文件"""
    with Image.open(uploaded_file) as image:
        return pytesseract.image_to_string(image)

def handle_csv(uploaded_file):
    """处理CSV文件"""
    text = ''
    # 使用文本模式读取文件，防止二进制模式下的编码错误
    with io.TextIOWrapper(uploaded_file, encoding='utf-8') as file:
        csv_reader = csv.reader(file)
        for row in csv_reader:
            text += '\t'.join(row) + '\n'
    return text

def handle_uploaded_file(uploaded_file):
    """处理上传的文件并返回文本内容。"""
    file_handlers = {
        "text/plain": handle_text,
        "application/pdf": handle_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": handle_docx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": handle_excel,
        "application/vnd.ms-excel": handle_excel,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": handle_ppt,
        "application/vnd.ms-powerpoint": handle_ppt,
        "image/jpeg": handle_image,
        "image/png": handle_image,
        "text/csv": handle_csv,  # CSV 文件的一个常见 MIME 类型
        "application/csv": handle_csv,  # 另一个可能的 MIME 类型
        "application/vnd.ms-excel": handle_csv  # CSV 文件可能被错误地识别为这个类型
    }

    try:
        handler = file_handlers.get(uploaded_file.type)
        if handler:
            return handler(uploaded_file)
        else:
            print(f"未识别的文件类型: {uploaded_file.type}")  # 打印未识别的文件类型
            raise ValueError("不支持的文件类型")
    except Exception as specific_exception:
        return f"处理文件时出错: {specific_exception}"


